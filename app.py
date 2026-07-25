from datetime import datetime, timezone
from functools import wraps
import json
import os
import re
import sqlite3

import numpy as np
import PyPDF2
import docx
from pptx import Presentation
import openpyxl
from flask import Flask, abort, jsonify, redirect, render_template, request, send_from_directory, session
from google import genai
from dotenv import load_dotenv

load_dotenv()

from keybert import KeyBERT
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from werkzeug.security import check_password_hash, generate_password_hash

app = Flask(__name__)
app.secret_key = "knowledge_portal_secret"
UPLOAD_FOLDER = "uploads"
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["MAX_CONTENT_LENGTH"] = 32 * 1024 * 1024  # 32 MB upload limit

os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# ─── SENTENCE TRANSFORMER ────────────────────────────────────────────────────
try:
    model = SentenceTransformer("all-MiniLM-L6-v2")
    kw_model = KeyBERT(model)
except Exception:
    try:
        model = SentenceTransformer("all-MiniLM-L6-v2", local_files_only=True)
        kw_model = KeyBERT(model)
    except Exception:
        model = None
        kw_model = None

# ─── CHROMADB SETUP ──────────────────────────────────────────────────────────
CHROMA_PATH = os.path.join(os.path.dirname(__file__), "chroma_db")
chroma_client = None
chroma_collection = None

try:
    import chromadb
    from chromadb.config import Settings

    chroma_client = chromadb.PersistentClient(
        path=CHROMA_PATH,
        settings=Settings(anonymized_telemetry=False),
    )
    chroma_collection = chroma_client.get_or_create_collection(
        name="knowledge_chunks",
        metadata={"hnsw:space": "cosine"},
    )
    print(f"[ChromaDB] Connected — collection has {chroma_collection.count()} vectors.")
except Exception as _chroma_err:
    print(f"[ChromaDB] Unavailable ({_chroma_err}). Falling back to in-memory cosine search.")
    chroma_client = None
    chroma_collection = None

# ─── LANGCHAIN SETUP ─────────────────────────────────────────────────────────
langchain_available = False
try:
    from langchain_google_genai import ChatGoogleGenerativeAI
    from langchain_core.prompts import PromptTemplate
    from langchain_core.output_parsers import StrOutputParser
    langchain_available = True
    print("[LangChain] Loaded successfully — quiz pipeline ready (LCEL enabled).")
except Exception as _lc_err:
    print(f"[LangChain] Unavailable ({_lc_err}). Falling back to direct REST API.")

# ─── CONSTANTS ────────────────────────────────────────────────────────────────
VALID_ROLES = ("student", "faculty", "administrator", "researcher")
ROLE_LABELS = {
    "student": "Student",
    "faculty": "Faculty",
    "administrator": "Administrator",
    "researcher": "Researcher",
}
AUDIENCE_OPTIONS = [
    ("all", "All Stakeholders"),
    ("students", "Students"),
    ("faculty", "Faculty"),
    ("researchers", "Researchers"),
    ("administrators", "Administrators"),
]


# ─────────────────────────────────────────────────────────────────────────────
# DATABASE
# ─────────────────────────────────────────────────────────────────────────────
def get_db():
    conn = sqlite3.connect("database.db")
    return conn


def ensure_column(cur, table_name, column_name, definition):
    cur.execute(f"PRAGMA table_info({table_name})")
    existing = {row[1] for row in cur.fetchall()}
    if column_name not in existing:
        cur.execute(f"ALTER TABLE {table_name} ADD COLUMN {column_name} {definition}")


def init_db():
    conn = get_db()
    cur = conn.cursor()

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS users(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT UNIQUE,
            password TEXT,
            role TEXT
        )
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS knowledge(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document TEXT,
            domain TEXT,
            text TEXT,
            embedding TEXT,
            audience TEXT DEFAULT 'all',
            keywords TEXT DEFAULT '',
            uploaded_by TEXT DEFAULT '',
            created_at TEXT DEFAULT '',
            chroma_id TEXT DEFAULT ''
        )
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS domains(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            name TEXT UNIQUE
        )
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS activity_log(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT,
            action TEXT,
            details TEXT,
            created_at TEXT
        )
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS quizzes(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            document TEXT,
            questions_json TEXT,
            created_at TEXT
        )
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS search_analytics(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT,
            query TEXT,
            results_count INTEGER,
            created_at TEXT
        )
        """
    )

    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS document_access_analytics(
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            username TEXT,
            document TEXT,
            access_type TEXT,
            created_at TEXT
        )
        """
    )

    # Migration support for old DBs.
    ensure_column(cur, "knowledge", "audience", "TEXT DEFAULT 'all'")
    ensure_column(cur, "knowledge", "keywords", "TEXT DEFAULT ''")
    ensure_column(cur, "knowledge", "uploaded_by", "TEXT DEFAULT ''")
    ensure_column(cur, "knowledge", "created_at", "TEXT DEFAULT ''")
    ensure_column(cur, "knowledge", "chroma_id", "TEXT DEFAULT ''")

    cur.execute("SELECT COUNT(*) FROM domains")
    if cur.fetchone()[0] == 0:
        cur.executemany(
            "INSERT INTO domains(name) VALUES (?)",
            [("Academics",), ("Administration",), ("Research",), ("General",)],
        )

    conn.commit()
    conn.close()


init_db()


# ─── CHROMADB MIGRATION: seed from SQLite on startup ─────────────────────────
def migrate_to_chroma():
    """If ChromaDB collection is empty but SQLite has chunks, re-index them."""
    if chroma_collection is None:
        return
    if chroma_collection.count() > 0:
        return  # already populated

    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        "SELECT id, document, domain, text, embedding, audience, keywords FROM knowledge ORDER BY id"
    )
    rows = cur.fetchall()
    conn.close()

    if not rows:
        return

    print(f"[ChromaDB] Migrating {len(rows)} existing chunks from SQLite …")
    batch_ids, batch_docs, batch_metas, batch_embeds = [], [], [], []

    for row_id, document, domain, text, embedding_json, audience, keywords in rows:
        chroma_id = f"chunk_{row_id}"
        try:
            vector = json.loads(embedding_json)
        except Exception:
            vector = encode_texts([text])[0].tolist()

        batch_ids.append(chroma_id)
        batch_docs.append(text)  # ChromaDB stores full doc for exact query support
        batch_metas.append({
            "sqlite_id": row_id,
            "document": document,
            "domain": domain or "",
            "audience": audience or "all",
            "keywords": keywords or "",
        })
        batch_embeds.append(vector)

        # Upsert in batches of 200
        if len(batch_ids) >= 200:
            chroma_collection.upsert(
                ids=batch_ids,
                documents=batch_docs,
                metadatas=batch_metas,
                embeddings=batch_embeds,
            )
            batch_ids, batch_docs, batch_metas, batch_embeds = [], [], [], []

    if batch_ids:
        chroma_collection.upsert(
            ids=batch_ids,
            documents=batch_docs,
            metadatas=batch_metas,
            embeddings=batch_embeds,
        )

    print(f"[ChromaDB] Migration complete — {chroma_collection.count()} vectors indexed.")


# ─── HELPERS ──────────────────────────────────────────────────────────────────
def normalize_role(role):
    role = (role or "").strip().lower()
    legacy_map = {"admin": "administrator", "user": "student"}
    return legacy_map.get(role, role)


def role_label(role):
    return ROLE_LABELS.get(role, "Stakeholder")


def allowed_audiences_for_role(role):
    role = normalize_role(role)
    mapping = {
        "administrator": {"all", "students", "faculty", "researchers", "administrators"},
        "faculty": {"all", "students", "faculty", "researchers"},
        "researcher": {"all", "researchers"},
        "student": {"all", "students"},
    }
    return mapping.get(role, {"all"})


def can_upload(role):
    return normalize_role(role) in {"faculty", "administrator", "researcher"}


def can_manage_domains(role):
    return normalize_role(role) == "administrator"


def log_activity(username, action, details):
    try:
        conn = get_db()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO activity_log(username,action,details,created_at) VALUES (?,?,?,?)",
            (username, action, details, datetime.now(timezone.utc).isoformat(timespec="seconds")),
        )
        conn.commit()
        conn.close()
    except Exception:
        pass


def log_document_access(username, document, access_type):
    try:
        conn = get_db()
        cur = conn.cursor()
        cur.execute(
            "INSERT INTO document_access_analytics(username, document, access_type, created_at) VALUES (?,?,?,?)",
            (username, document, access_type, datetime.now(timezone.utc).isoformat(timespec="seconds")),
        )
        conn.commit()
        conn.close()
    except Exception:
        pass


def login_required(route_fn):
    @wraps(route_fn)
    def wrapper(*args, **kwargs):
        if not session.get("username"):
            return redirect("/")
        return route_fn(*args, **kwargs)

    return wrapper


def roles_required(*roles):
    allowed_roles = {normalize_role(role) for role in roles}

    def decorator(route_fn):
        @wraps(route_fn)
        def wrapper(*args, **kwargs):
            current_role = normalize_role(session.get("role"))
            if not session.get("username"):
                return redirect("/")
            if current_role not in allowed_roles:
                return abort(403)
            return route_fn(*args, **kwargs)

        return wrapper

    return decorator


# ─── TEXT EXTRACTORS ──────────────────────────────────────────────────────────
def extract_text_from_pdf(filepath):
    text = ""
    with open(filepath, "rb") as file_handle:
        reader = PyPDF2.PdfReader(file_handle)
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
    return text


def extract_text_from_docx(filepath):
    doc = docx.Document(filepath)
    full_text = []
    for paragraph in doc.paragraphs:
        if paragraph.text.strip():
            full_text.append(paragraph.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if row_text:
                full_text.append(" | ".join(row_text))
    return "\n".join(full_text)


def extract_text_from_pptx(filepath):
    prs = Presentation(filepath)
    text_runs = []
    for i, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            slide_text.append(run.text.strip())
        if slide_text:
            text_runs.append(f"[Slide {i+1}] " + " ".join(slide_text))
    return "\n".join(text_runs)


def extract_text_from_xlsx(filepath):
    wb = openpyxl.load_workbook(filepath, data_only=True)
    text_runs = []
    for sheet_name in wb.sheetnames:
        sheet = wb[sheet_name]
        sheet_text = []
        for row in sheet.iter_rows(values_only=True):
            row_vals = [str(val).strip() for val in row if val is not None]
            if row_vals:
                sheet_text.append(" | ".join(row_vals))
        if sheet_text:
            text_runs.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_text))
    return "\n\n".join(text_runs)


def extract_text_from_txt(filepath):
    try:
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
    except UnicodeDecodeError:
        with open(filepath, "r", encoding="latin-1") as f:
            return f.read()


def split_into_chunks(text, chunk_size=180, overlap=40):
    words = text.split()
    if not words:
        return []

    chunks = []
    step = max(1, chunk_size - overlap)
    for start in range(0, len(words), step):
        chunk_words = words[start : start + chunk_size]
        if not chunk_words:
            continue
        chunk = " ".join(chunk_words).strip()
        if len(chunk) > 50:
            chunks.append(chunk)
        if start + chunk_size >= len(words):
            break
    return chunks


def hashed_embedding(text, dim=384):
    """Deterministic hash-based fallback embedding when no ML model is available."""
    vector = np.zeros(dim, dtype=np.float32)
    tokens = re.findall(r"[a-zA-Z0-9]{2,}", text.lower())
    for token in tokens:
        vector[hash(token) % dim] += 1.0
    norm = np.linalg.norm(vector)
    if norm > 0:
        vector = vector / norm
    return vector


def encode_texts(texts):
    if not texts:
        return np.array([], dtype=np.float32)
    if model is not None:
        return model.encode(texts)
    return np.array([hashed_embedding(text) for text in texts], dtype=np.float32)


def auto_keywords(text):
    sample = " ".join(text.split()[:500])
    if not sample:
        return ""

    if kw_model is not None:
        try:
            extracted = kw_model.extract_keywords(
                sample,
                keyphrase_ngram_range=(1, 2),
                stop_words="english",
                top_n=5,
            )
            return ", ".join([kw for kw, _ in extracted])
        except Exception:
            pass

    tokens = re.findall(r"[a-zA-Z]{4,}", sample.lower())
    unique = []
    for token in tokens:
        if token not in unique:
            unique.append(token)
        if len(unique) == 5:
            break
    return ", ".join(unique)


def is_document_accessible(filename, role):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT audience FROM knowledge WHERE document=?", (filename,))
    audiences = [row[0] or "all" for row in cur.fetchall()]
    conn.close()
    if not audiences:
        return False
    allowed = allowed_audiences_for_role(role)
    return any(audience in allowed for audience in audiences)


def get_dashboard_stats(role):
    allowed = tuple(allowed_audiences_for_role(role))
    placeholders = ",".join(["?"] * len(allowed))

    conn = get_db()
    cur = conn.cursor()

    cur.execute(
        f"SELECT COUNT(DISTINCT document) FROM knowledge WHERE audience IN ({placeholders})",
        allowed,
    )
    doc_count = cur.fetchone()[0]

    cur.execute("SELECT COUNT(*) FROM domains")
    domain_count = cur.fetchone()[0]

    cur.execute(
        f"""
        SELECT document, domain, MAX(created_at)
        FROM knowledge
        WHERE audience IN ({placeholders})
        GROUP BY document, domain
        ORDER BY MAX(created_at) DESC
        LIMIT 5
        """,
        allowed,
    )
    recent_docs = cur.fetchall()

    cur.execute("SELECT COUNT(*) FROM users")
    user_count = cur.fetchone()[0]

    top_queries = []
    zero_queries = []
    top_docs = []

    if role == "administrator":
        try:
            cur.execute(
                """
                SELECT query, COUNT(*) as cnt
                FROM search_analytics
                GROUP BY query
                ORDER BY cnt DESC
                LIMIT 5
                """
            )
            top_queries = cur.fetchall()

            cur.execute(
                """
                SELECT query, COUNT(*) as cnt
                FROM search_analytics
                WHERE results_count = 0
                GROUP BY query
                ORDER BY cnt DESC
                LIMIT 5
                """
            )
            zero_queries = cur.fetchall()

            cur.execute(
                """
                SELECT document, COUNT(*) as cnt
                FROM document_access_analytics
                GROUP BY document
                ORDER BY cnt DESC
                LIMIT 5
                """
            )
            top_docs = cur.fetchall()
        except Exception:
            pass

    conn.close()

    chroma_count = chroma_collection.count() if chroma_collection else 0

    return {
        "doc_count": doc_count,
        "domain_count": domain_count,
        "recent_docs": recent_docs,
        "user_count": user_count,
        "top_queries": top_queries,
        "zero_queries": zero_queries,
        "top_docs": top_docs,
        "chroma_vectors": chroma_count,
        "langchain_enabled": langchain_available,
    }


# ─── CHROMADB SEARCH ──────────────────────────────────────────────────────────
def search_with_chroma(query_text, allowed_audiences, domain_filter=None, n_results=10):
    """
    Semantic search using ChromaDB vector store.
    Returns list of dicts: {document, domain, text, audience, keywords, score}
    """
    if chroma_collection is None or chroma_collection.count() == 0:
        return None  # signal caller to fall back to SQLite

    query_embedding = encode_texts([query_text])[0].tolist()

    where_filter = {"audience": {"$in": list(allowed_audiences)}}
    if domain_filter and domain_filter != "all":
        where_filter = {
            "$and": [
                {"audience": {"$in": list(allowed_audiences)}},
                {"domain": {"$eq": domain_filter}},
            ]
        }

    try:
        results = chroma_collection.query(
            query_embeddings=[query_embedding],
            n_results=min(n_results, chroma_collection.count()),
            where=where_filter,
            include=["documents", "metadatas", "distances"],
        )
    except Exception as e:
        print(f"[ChromaDB] Query error: {e}")
        return None

    hits = []
    ids = results.get("ids", [[]])[0]
    docs = results.get("documents", [[]])[0]
    metas = results.get("metadatas", [[]])[0]
    distances = results.get("distances", [[]])[0]

    for i, (meta, distance) in enumerate(zip(metas, distances)):
        # ChromaDB cosine distance → similarity score (0–100)
        similarity = max(0.0, 1.0 - distance)
        hits.append({
            "document": meta.get("document", ""),
            "domain": meta.get("domain", ""),
            "text": docs[i] if i < len(docs) else "",
            "audience": meta.get("audience", "all"),
            "keywords": meta.get("keywords", ""),
            "score": int(similarity * 100),
            "sqlite_id": meta.get("sqlite_id"),
        })

    return hits


# ─── LANGCHAIN QUIZ GENERATION ────────────────────────────────────────────────
QUIZ_PROMPT_TEMPLATE = """You are an expert academic educator. Based on the following source document text, generate exactly {num_questions} Multiple Choice Questions (MCQs) to test a student's comprehension.

Source Text:
{source_text}

Provide the output ONLY as a raw JSON array of objects. Do not wrap the JSON in markdown code blocks. Each object must have exactly these keys:
- "question": The question string.
- "options": A list of exactly 4 strings for options A, B, C, D in order.
- "answer": A single uppercase character: "A", "B", "C", or "D".
- "explanation": A detailed explanation of why the correct option is right.

Ensure options are clear, grammatically correct, and only one option is correct."""


def generate_quiz_with_langchain(api_key: str, source_text: str, num_questions: int) -> list:
    """
    Generate quiz questions using LangChain LLMChain with ChatGoogleGenerativeAI.
    Returns a list of question dicts.
    Raises RuntimeError on failure.
    """
    if not langchain_available:
        raise ImportError("LangChain not available")

    llm = ChatGoogleGenerativeAI(
        model="gemini-1.5-flash",
        google_api_key=api_key,
        temperature=0.4,
        max_output_tokens=8192,
    )

    prompt = PromptTemplate(
        input_variables=["num_questions", "source_text"],
        template=QUIZ_PROMPT_TEMPLATE,
    )

    chain = prompt | llm | StrOutputParser()
    response_text = chain.invoke({"num_questions": num_questions, "source_text": source_text})
    response_text = (response_text or "").strip()
    if not response_text:
        raise ValueError("LangChain chain returned empty response")

    # Strip markdown fences if present
    if response_text.startswith("```"):
        response_text = re.sub(r"^```[a-zA-Z]*\n?", "", response_text)
        response_text = re.sub(r"\n?```$", "", response_text)
    response_text = response_text.strip()

    questions = json.loads(response_text)
    if not isinstance(questions, list) or len(questions) == 0:
        raise ValueError("LangChain returned no valid questions")
    return questions


def generate_quiz_with_rest(api_key: str, source_text: str, num_questions: int) -> list:
    """
    Fallback: generate quiz questions using direct Gemini REST API.
    Tries multiple model IDs in order.
    """
    import urllib.request
    import urllib.error

    prompt = QUIZ_PROMPT_TEMPLATE.format(
        num_questions=num_questions,
        source_text=source_text,
    )

    MODELS_TO_TRY = [
        "gemini-flash-latest",
        "gemini-pro-latest",
        "gemini-2.0-flash",
        "gemini-2.0-flash-lite",
        "gemini-2.5-flash",
        "gemini-2.5-flash-lite",
        "gemini-2.5-pro",
    ]

    response_text = None
    last_error = None

    for model_id in MODELS_TO_TRY:
        url = (
            f"https://generativelanguage.googleapis.com/v1beta/models/"
            f"{model_id}:generateContent?key={api_key}"
        )
        payload = json.dumps({
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.4, "maxOutputTokens": 8192},
        }).encode("utf-8")

        req = urllib.request.Request(
            url,
            data=payload,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        try:
            with urllib.request.urlopen(req, timeout=120) as resp:
                result = json.loads(resp.read().decode("utf-8"))
                response_text = result["candidates"][0]["content"]["parts"][0]["text"].strip()
                break
        except urllib.error.HTTPError as http_err:
            err_body = http_err.read().decode("utf-8", errors="ignore")
            last_error = f"HTTP {http_err.code} from {model_id}: {err_body[:300]}"
            if http_err.code in (404, 400, 429):
                continue
            raise RuntimeError(last_error)
        except Exception as e:
            last_error = str(e)
            continue

    if response_text is None:
        raise RuntimeError(
            f"All Gemini models are unavailable or rate-limited. Last error: {last_error}"
        )

    if response_text.startswith("```"):
        response_text = re.sub(r"^```[a-zA-Z]*\n?", "", response_text)
        response_text = re.sub(r"\n?```$", "", response_text)
    response_text = response_text.strip()

    questions = json.loads(response_text)
    if not isinstance(questions, list) or len(questions) == 0:
        raise ValueError("LLM did not return a valid list of questions.")
    return questions


# ─────────────────────────────────────────────────────────────────────────────
# AUTH ROUTES
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/")
def home():
    if session.get("username"):
        return redirect("/dashboard")
    return render_template("index.html")


@app.route("/register", methods=["GET", "POST"])
def register():
    if request.method == "POST":
        username = request.form["username"].strip()
        password = request.form["password"]
        role = normalize_role(request.form["role"])

        if role not in VALID_ROLES:
            return render_template("register.html", error="Please choose a valid role.", roles=VALID_ROLES)

        conn = get_db()
        cur = conn.cursor()
        cur.execute("SELECT id FROM users WHERE username=?", (username,))
        existing = cur.fetchone()
        if existing:
            conn.close()
            return render_template("register.html", error="Username already exists.", roles=VALID_ROLES)

        cur.execute(
            "INSERT INTO users(username,password,role) VALUES (?,?,?)",
            (username, generate_password_hash(password), role),
        )
        conn.commit()
        conn.close()
        log_activity(username, "register", f"New {role} account created")
        return redirect("/")

    return render_template("register.html", roles=VALID_ROLES)


@app.route("/login", methods=["POST"])
def login():
    username = request.form["username"].strip()
    password = request.form["password"]

    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT id,password,role FROM users WHERE username=?", (username,))
    user = cur.fetchone()

    if not user:
        conn.close()
        return render_template("index.html", error="Invalid username or password.")

    user_id, stored_password, role = user
    role = normalize_role(role)
    password_ok = False

    try:
        password_ok = check_password_hash(stored_password, password)
    except Exception:
        password_ok = False

    # Legacy plaintext record fallback.
    if not password_ok and ":" not in (stored_password or ""):
        password_ok = stored_password == password
        if password_ok:
            cur.execute(
                "UPDATE users SET password=?, role=? WHERE id=?",
                (generate_password_hash(password), role, user_id),
            )
            conn.commit()

    conn.close()

    if not password_ok:
        return render_template("index.html", error="Invalid username or password.")

    session["user_id"] = user_id
    session["username"] = username
    session["role"] = role
    log_activity(username, "login", f"{role} logged in")
    return redirect("/dashboard")


@app.route("/logout")
def logout():
    if session.get("username"):
        log_activity(session["username"], "logout", "User logged out")
    session.clear()
    return redirect("/")


# ─────────────────────────────────────────────────────────────────────────────
# DASHBOARD
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/dashboard")
@login_required
def dashboard():
    role = normalize_role(session.get("role"))
    stats = get_dashboard_stats(role)
    return render_template(
        "dashboard.html",
        role=role,
        role_label=role_label(role),
        username=session.get("username"),
        can_upload_docs=can_upload(role),
        can_manage_domain=can_manage_domains(role),
        stats=stats,
    )


@app.route("/admin")
@roles_required("administrator")
def admin_dashboard():
    return redirect("/dashboard")


@app.route("/user")
@login_required
def user():
    return redirect("/dashboard")


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENT UPLOAD
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/upload", methods=["GET", "POST"])
@roles_required("faculty", "administrator", "researcher")
def upload():
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT name FROM domains ORDER BY name")
    domains = [row[0] for row in cur.fetchall()]
    conn.close()

    if request.method == "POST":
        files = request.files.getlist("files") or request.files.getlist("pdf")
        domain = request.form["domain"]
        audience = request.form.get("audience", "all")
        manual_keywords = request.form.get("keywords", "").strip()
        username = session.get("username", "unknown")

        if audience not in {item[0] for item in AUDIENCE_OPTIONS}:
            return render_template(
                "upload.html",
                domains=domains,
                audience_options=AUDIENCE_OPTIONS,
                error="Invalid audience selected.",
            )

        saved_docs = 0
        conn = get_db()
        cur = conn.cursor()
        chroma_batch_ids, chroma_batch_docs, chroma_batch_metas, chroma_batch_embeds = [], [], [], []

        for file in files:
            if not file.filename:
                continue
            filename_lower = file.filename.lower()
            if not filename_lower.endswith((".pdf", ".docx", ".pptx", ".xlsx", ".txt")):
                continue

            filename = os.path.basename(file.filename)
            filepath = os.path.join(app.config["UPLOAD_FOLDER"], filename)
            file.save(filepath)

            if filename_lower.endswith(".pdf"):
                text = extract_text_from_pdf(filepath)
            elif filename_lower.endswith(".docx"):
                text = extract_text_from_docx(filepath)
            elif filename_lower.endswith(".pptx"):
                text = extract_text_from_pptx(filepath)
            elif filename_lower.endswith(".xlsx"):
                text = extract_text_from_xlsx(filepath)
            elif filename_lower.endswith(".txt"):
                text = extract_text_from_txt(filepath)
            else:
                text = ""

            chunks = split_into_chunks(text)
            if not chunks:
                continue

            doc_keywords = manual_keywords or auto_keywords(text)
            embeddings = encode_texts(chunks)
            created_at = datetime.now(timezone.utc).isoformat(timespec="seconds")

            for idx, chunk in enumerate(chunks):
                embedding_list = embeddings[idx].tolist()
                cur.execute(
                    """
                    INSERT INTO knowledge(document,domain,text,embedding,audience,keywords,uploaded_by,created_at)
                    VALUES (?,?,?,?,?,?,?,?)
                    """,
                    (
                        filename,
                        domain,
                        chunk,
                        json.dumps(embedding_list),
                        audience,
                        doc_keywords,
                        username,
                        created_at,
                    ),
                )
                sqlite_id = cur.lastrowid
                chroma_id = f"chunk_{sqlite_id}"

                # Update the chroma_id back into SQLite
                cur.execute("UPDATE knowledge SET chroma_id=? WHERE id=?", (chroma_id, sqlite_id))

                # Stage for ChromaDB batch upsert
                chroma_batch_ids.append(chroma_id)
                chroma_batch_docs.append(chunk)
                chroma_batch_metas.append({
                    "sqlite_id": sqlite_id,
                    "document": filename,
                    "domain": domain,
                    "audience": audience,
                    "keywords": doc_keywords,
                })
                chroma_batch_embeds.append(embedding_list)

            saved_docs += 1
            log_activity(username, "upload", f"Uploaded {filename} in {domain}")

        conn.commit()
        conn.close()

        # Batch upsert into ChromaDB
        if chroma_collection is not None and chroma_batch_ids:
            try:
                BATCH_SIZE = 200
                for i in range(0, len(chroma_batch_ids), BATCH_SIZE):
                    chroma_collection.upsert(
                        ids=chroma_batch_ids[i:i+BATCH_SIZE],
                        documents=chroma_batch_docs[i:i+BATCH_SIZE],
                        metadatas=chroma_batch_metas[i:i+BATCH_SIZE],
                        embeddings=chroma_batch_embeds[i:i+BATCH_SIZE],
                    )
                print(f"[ChromaDB] Indexed {len(chroma_batch_ids)} new vectors. Total: {chroma_collection.count()}")
            except Exception as e:
                print(f"[ChromaDB] Upsert error (non-fatal): {e}")

        message = f"{saved_docs} document(s) uploaded and indexed successfully."
        return render_template(
            "upload.html",
            domains=domains,
            audience_options=AUDIENCE_OPTIONS,
            success=message,
        )

    return render_template("upload.html", domains=domains, audience_options=AUDIENCE_OPTIONS)


# ─────────────────────────────────────────────────────────────────────────────
# SEMANTIC SEARCH
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/search", methods=["GET", "POST"])
@login_required
def search():
    role = normalize_role(session.get("role"))
    is_admin = role == "administrator"
    allowed_audiences = allowed_audiences_for_role(role)
    audience_choices = [item for item in AUDIENCE_OPTIONS if item[0] in allowed_audiences or item[0] == "all"]

    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT name FROM domains ORDER BY name")
    domains = [row[0] for row in cur.fetchall()]
    conn.close()

    results = []
    selected_domain = "all"
    selected_audience = "all"
    search_query = ""
    search_mode = "hybrid"
    search_engine_used = "chromadb" if chroma_collection and chroma_collection.count() > 0 else "sqlite"

    if request.method == "POST":
        search_query = request.form["query"].strip()
        selected_domain = request.form.get("domain", "all")
        selected_audience = request.form.get("audience", "all") if is_admin else "all"
        search_mode = request.form.get("search_mode", "hybrid")

        if search_query:
            if search_mode == "exact":
                # ── EXACT MATCH RETRIEVAL (Bypass AI vector matching) ───────
                search_engine_used = "sqlite_exact"
                allowed_tuple = tuple(allowed_audiences)
                placeholders = ",".join(["?"] * len(allowed_tuple))
                sql = (
                    "SELECT document, domain, text, audience, keywords "
                    f"FROM knowledge WHERE audience IN ({placeholders}) AND "
                    "(text LIKE ? OR document LIKE ? OR keywords LIKE ?)"
                )
                params = list(allowed_tuple) + [f"%{search_query}%", f"%{search_query}%", f"%{search_query}%"]

                if selected_domain != "all":
                    sql += " AND domain=?"
                    params.append(selected_domain)
                if is_admin and selected_audience != "all" and selected_audience in allowed_audiences:
                    sql += " AND audience=?"
                    params.append(selected_audience)

                conn = get_db()
                cur = conn.cursor()
                cur.execute(sql, tuple(params))
                rows = cur.fetchall()
                conn.close()

                ranked = []
                q_lower = search_query.lower()
                for doc, domain, text, audience, keywords in rows:
                    score = 0
                    if q_lower in doc.lower():
                        score += 50
                    if q_lower in (keywords or "").lower():
                        score += 30
                    if q_lower in text.lower():
                        score += 20
                    
                    # Frequency bonus for exact matches in text
                    freq = text.lower().count(q_lower)
                    score += min(10, freq * 2)
                    score = min(100, score)
                    ranked.append((doc, domain, text, audience, keywords, score))

                ranked.sort(key=lambda x: x[5], reverse=True)

                seen_documents = set()
                for doc, domain, text, audience, keywords, score in ranked:
                    if doc in seen_documents:
                        continue
                    results.append({
                        "document": doc,
                        "domain": domain,
                        "audience": audience,
                        "keywords": keywords,
                        "snippet": text[:240],
                        "score": score,
                    })
                    seen_documents.add(doc)
                    if len(results) == 5:
                        break

            else:
                # ── HYBRID RETRIEVAL (Semantic vector search + keywords) ─────
                chroma_hits = search_with_chroma(
                    search_query,
                    allowed_audiences,
                    domain_filter=selected_domain,
                    n_results=10,
                )

                if chroma_hits is not None:
                    # ChromaDB returned results — deduplicate by document
                    seen_documents = set()
                    query_terms = set(re.findall(r"[a-zA-Z]{3,}", search_query.lower()))

                    for hit in chroma_hits:
                        doc = hit["document"]
                        if doc in seen_documents:
                            continue
                        if hit["score"] < 5:
                            continue

                        # Optionally boost with keyword overlap
                        keyword_source = f"{hit['keywords']} {hit['text']}".lower()
                        keyword_terms = set(re.findall(r"[a-zA-Z]{3,}", keyword_source))
                        keyword_overlap = 0.0
                        if query_terms:
                            keyword_overlap = len(query_terms.intersection(keyword_terms)) / len(query_terms)

                        final_score = (0.85 * hit["score"]) + (0.15 * keyword_overlap * 100)

                        # Exact phrase boost (+15 points if the exact query appears in the chunk text or title)
                        if search_query.lower() in hit["text"].lower() or search_query.lower() in doc.lower():
                            final_score += 15

                        final_score = int(min(100, final_score))

                        results.append({
                            "document": doc,
                            "domain": hit["domain"],
                            "audience": hit["audience"],
                            "keywords": hit["keywords"],
                            "snippet": hit["text"][:240],
                            "score": final_score,
                        })
                        seen_documents.add(doc)
                        if len(results) == 5:
                            break

                else:
                    # ── FALLBACK: SQLite in-memory cosine search ─────────────
                    search_engine_used = "sqlite_fallback"
                    allowed_tuple = tuple(allowed_audiences)
                    placeholders = ",".join(["?"] * len(allowed_tuple))
                    sql = (
                        "SELECT document,domain,text,embedding,audience,keywords "
                        f"FROM knowledge WHERE audience IN ({placeholders})"
                    )
                    params = list(allowed_tuple)

                    if selected_domain != "all":
                        sql += " AND domain=?"
                        params.append(selected_domain)
                    if is_admin and selected_audience != "all" and selected_audience in allowed_audiences:
                        sql += " AND audience=?"
                        params.append(selected_audience)

                    conn = get_db()
                    cur = conn.cursor()
                    cur.execute(sql, tuple(params))
                    rows = cur.fetchall()
                    conn.close()

                    if rows:
                        valid_rows = []
                        embedding_vectors = []

                        for doc, domain, text, embedding_json, audience, keywords in rows:
                            try:
                                vector = np.array(json.loads(embedding_json), dtype=np.float32)
                            except Exception:
                                vector = encode_texts([text])[0]

                            valid_rows.append((doc, domain, text, audience, keywords or ""))
                            embedding_vectors.append(vector)

                        embeddings_matrix = np.array(embedding_vectors, dtype=np.float32)
                        query_embedding = encode_texts([search_query])
                        semantic_scores = cosine_similarity(query_embedding, embeddings_matrix)[0]
                        query_terms = set(re.findall(r"[a-zA-Z]{3,}", search_query.lower()))

                        ranked = []
                        for row_data, semantic in zip(valid_rows, semantic_scores):
                            _, _, text, _, keywords = row_data
                            keyword_source = f"{keywords} {text}".lower()
                            keyword_terms = set(re.findall(r"[a-zA-Z]{3,}", keyword_source))
                            keyword_overlap = 0.0
                            if query_terms:
                                keyword_overlap = len(query_terms.intersection(keyword_terms)) / len(query_terms)
                            final_score = (0.85 * float(semantic)) + (0.15 * keyword_overlap)

                            # Exact phrase boost
                            if search_query.lower() in text.lower() or search_query.lower() in row_data[0].lower():
                                final_score += 0.15

                            ranked.append((row_data, final_score))

                        ranked.sort(key=lambda item: item[1], reverse=True)

                        seen_documents = set()
                        for (doc, domain, text, audience, keywords), score in ranked:
                            if doc in seen_documents:
                                continue
                            if score < 0.05:
                                continue
                            results.append({
                                "document": doc,
                                "domain": domain,
                                "audience": audience,
                                "keywords": keywords,
                                "snippet": text[:240],
                                "score": int(min(100, score * 100)),
                            })
                            seen_documents.add(doc)
                            if len(results) == 5:
                                break

        # Log search analytics
        try:
            conn_log = get_db()
            cur_log = conn_log.cursor()
            cur_log.execute(
                "INSERT INTO search_analytics(username, query, results_count, created_at) VALUES (?,?,?,?)",
                (
                    session.get("username", "anonymous"),
                    search_query,
                    len(results),
                    datetime.now(timezone.utc).isoformat(timespec="seconds"),
                ),
            )
            conn_log.commit()
            conn_log.close()
        except Exception:
            pass

    return render_template(
        "search.html",
        results=results,
        role=role,
        domains=domains,
        selected_domain=selected_domain,
        selected_audience=selected_audience,
        audience_options=audience_choices,
        can_filter_audience=is_admin,
        search_query=search_query,
        search_engine=search_engine_used,
        search_mode=search_mode,
    )


# ─────────────────────────────────────────────────────────────────────────────
# REST API — /api/search  (JSON endpoint)
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/api/search", methods=["GET", "POST"])
@login_required
def api_search():
    """
    JSON API endpoint for semantic search.
    Accepts query param 'q' (GET or POST JSON body).
    Returns: {"results": [...], "engine": "chromadb"|"sqlite_fallback", "count": N}
    """
    if request.method == "POST":
        data = request.get_json(silent=True) or {}
        query = data.get("q", "").strip()
        domain_filter = data.get("domain", "all")
    else:
        query = request.args.get("q", "").strip()
        domain_filter = request.args.get("domain", "all")

    if not query:
        return jsonify({"error": "Missing query parameter 'q'"}), 400

    role = normalize_role(session.get("role"))
    allowed_audiences = allowed_audiences_for_role(role)
    engine_used = "chromadb"

    chroma_hits = search_with_chroma(query, allowed_audiences, domain_filter=domain_filter, n_results=5)

    if chroma_hits is not None:
        api_results = [
            {
                "document": h["document"],
                "domain": h["domain"],
                "score": h["score"],
                "snippet": h["text"][:200],
                "keywords": h["keywords"],
            }
            for h in chroma_hits
            if h["score"] >= 5
        ]
    else:
        engine_used = "sqlite_fallback"
        allowed_tuple = tuple(allowed_audiences)
        placeholders = ",".join(["?"] * len(allowed_tuple))
        conn = get_db()
        cur = conn.cursor()
        cur.execute(
            f"SELECT document,domain,text,embedding,keywords FROM knowledge WHERE audience IN ({placeholders})",
            allowed_tuple,
        )
        rows = cur.fetchall()
        conn.close()
        query_emb = encode_texts([query])
        hits = []
        for doc, domain, text, emb_json, keywords in rows:
            try:
                vec = np.array(json.loads(emb_json), dtype=np.float32)
            except Exception:
                vec = encode_texts([text])[0]
            sim = float(cosine_similarity(query_emb, [vec])[0][0])
            hits.append((doc, domain, text, keywords, sim))
        hits.sort(key=lambda x: x[4], reverse=True)
        seen = set()
        api_results = []
        for doc, domain, text, keywords, sim in hits:
            if doc in seen or sim < 0.05:
                continue
            api_results.append({"document": doc, "domain": domain, "score": int(sim * 100), "snippet": text[:200], "keywords": keywords})
            seen.add(doc)
            if len(api_results) == 5:
                break

    # Log
    try:
        conn_log = get_db()
        cur_log = conn_log.cursor()
        cur_log.execute(
            "INSERT INTO search_analytics(username, query, results_count, created_at) VALUES (?,?,?,?)",
            (session.get("username", "api"), query, len(api_results), datetime.now(timezone.utc).isoformat(timespec="seconds")),
        )
        conn_log.commit()
        conn_log.close()
    except Exception:
        pass

    return jsonify({"results": api_results, "engine": engine_used, "count": len(api_results)})


# ─────────────────────────────────────────────────────────────────────────────
# PREVIEW / DOWNLOAD
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/preview/<filename>")
@login_required
def preview(filename):
    role = normalize_role(session.get("role"))
    if not is_document_accessible(filename, role):
        return abort(403)
    log_document_access(session.get("username", "anonymous"), filename, "preview")
    return send_from_directory("uploads", filename)


@app.route("/download/<filename>")
@login_required
def download(filename):
    role = normalize_role(session.get("role"))
    if not is_document_accessible(filename, role):
        return abort(403)
    log_document_access(session.get("username", "anonymous"), filename, "download")
    return send_from_directory("uploads", filename, as_attachment=True)


# ─────────────────────────────────────────────────────────────────────────────
# DOCUMENTS LIBRARY
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/documents")
@login_required
def documents():
    role = normalize_role(session.get("role"))
    allowed = tuple(allowed_audiences_for_role(role))
    placeholders = ",".join(["?"] * len(allowed))

    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        f"""
        SELECT
            domain,
            document,
            audience,
            COALESCE(MAX(keywords), ''),
            COALESCE(MAX(created_at), '')
        FROM knowledge
        WHERE audience IN ({placeholders})
        GROUP BY domain, document, audience
        ORDER BY domain, document
        """,
        allowed,
    )
    rows = cur.fetchall()
    conn.close()

    data = {}
    categories = set()
    for domain, doc, audience, keywords, created_at in rows:
        categories.add(domain)
        if domain not in data:
            data[domain] = []
        data[domain].append(
            {
                "name": doc,
                "audience": audience,
                "keywords": keywords,
                "created_at": created_at,
            }
        )

    return render_template(
        "documents.html",
        data=data,
        categories=sorted(categories),
        role=role,
    )


# ─────────────────────────────────────────────────────────────────────────────
# DOMAINS
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/domains", methods=["GET", "POST"])
@roles_required("administrator")
def domains():
    message = None
    error = None
    conn = get_db()
    cur = conn.cursor()

    if request.method == "POST":
        name = request.form["name"].strip()
        if not name:
            error = "Domain name cannot be empty."
        else:
            try:
                cur.execute("INSERT INTO domains(name) VALUES (?)", (name,))
                conn.commit()
                message = "Domain added successfully."
                log_activity(session.get("username", "unknown"), "add_domain", name)
            except sqlite3.IntegrityError:
                error = "Domain already exists."

    cur.execute("SELECT name FROM domains ORDER BY name")
    rows = cur.fetchall()
    conn.close()

    return render_template("domains.html", domains=rows, message=message, error=error)


# ─────────────────────────────────────────────────────────────────────────────
# AI QUIZ GENERATOR  (LangChain-powered, REST fallback)
# ─────────────────────────────────────────────────────────────────────────────
@app.route("/quiz", methods=["GET"])
@login_required
def quiz_dashboard():
    role = normalize_role(session.get("role"))
    allowed = tuple(allowed_audiences_for_role(role))
    placeholders = ",".join(["?"] * len(allowed))

    conn = get_db()
    cur = conn.cursor()
    cur.execute(
        f"SELECT DISTINCT document, domain FROM knowledge WHERE audience IN ({placeholders}) ORDER BY document",
        allowed,
    )
    docs = cur.fetchall()

    cur.execute("SELECT id, document, created_at FROM quizzes ORDER BY id DESC")
    all_quizzes = cur.fetchall()
    conn.close()

    accessible_docs = {doc[0] for doc in docs}
    user_quizzes = [q for q in all_quizzes if q[1] in accessible_docs]

    has_env_key = bool(os.environ.get("GEMINI_API_KEY"))

    return render_template(
        "quiz_generator.html",
        docs=docs,
        quizzes=user_quizzes,
        has_env_key=has_env_key,
        langchain_enabled=langchain_available,
        role=role,
    )


@app.route("/quiz/generate", methods=["POST"])
@login_required
def quiz_generate():
    document = request.form.get("document", "").strip()
    num_questions = int(request.form.get("num_questions", "5"))
    custom_api_key = request.form.get("custom_api_key", "").strip()

    role = normalize_role(session.get("role"))
    if not is_document_accessible(document, role):
        return abort(403)

    api_key = custom_api_key or os.environ.get("GEMINI_API_KEY")

    def _render_quiz_error(error_msg):
        allowed = tuple(allowed_audiences_for_role(role))
        placeholders = ",".join(["?"] * len(allowed))
        conn = get_db()
        cur = conn.cursor()
        cur.execute(
            f"SELECT DISTINCT document, domain FROM knowledge WHERE audience IN ({placeholders}) ORDER BY document",
            allowed,
        )
        docs = cur.fetchall()
        cur.execute("SELECT id, document, created_at FROM quizzes ORDER BY id DESC")
        all_quizzes = cur.fetchall()
        conn.close()
        accessible_docs = {doc[0] for doc in docs}
        user_quizzes = [q for q in all_quizzes if q[1] in accessible_docs]
        return render_template(
            "quiz_generator.html",
            docs=docs,
            quizzes=user_quizzes,
            has_env_key=bool(os.environ.get("GEMINI_API_KEY")),
            langchain_enabled=langchain_available,
            role=role,
            error=error_msg,
        )

    if not api_key:
        return _render_quiz_error(
            "Gemini API Key is required. Please set GEMINI_API_KEY in your .env file or enter it below."
        )

    # Retrieve document chunks from SQLite
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT text FROM knowledge WHERE document=? ORDER BY id", (document,))
    chunks = [row[0] for row in cur.fetchall()]
    conn.close()

    if not chunks:
        return _render_quiz_error(f"No content found for '{document}'. Please re-upload the document.")

    full_text = "\n".join(chunks)
    words = full_text.split()
    if len(words) > 15000:
        full_text = " ".join(words[:15000])

    try:
        questions = None
        pipeline_used = "rest_api"

        # ── PRIMARY: LangChain LLMChain pipeline ─────────────────────────────
        if langchain_available:
            try:
                questions = generate_quiz_with_langchain(api_key, full_text, num_questions)
                pipeline_used = "langchain"
                print(f"[LangChain] Successfully generated {len(questions)} questions for '{document}'")
            except Exception as lc_err:
                print(f"[LangChain] Failed ({lc_err}), falling back to REST API …")
                questions = None

        # ── FALLBACK: Direct REST API ─────────────────────────────────────────
        if questions is None:
            questions = generate_quiz_with_rest(api_key, full_text, num_questions)
            pipeline_used = "rest_api"

        conn = get_db()
        cur = conn.cursor()
        created_at = datetime.now(timezone.utc).isoformat(timespec="seconds")
        cur.execute(
            "INSERT INTO quizzes (document, questions_json, created_at) VALUES (?, ?, ?)",
            (document, json.dumps(questions), created_at)
        )
        quiz_id = cur.lastrowid
        conn.commit()
        conn.close()

        log_activity(
            session.get("username", "unknown"),
            "generate_quiz",
            f"Generated quiz {quiz_id} for '{document}' via {pipeline_used} ({len(questions)} questions)",
        )
        return redirect(f"/quiz/view/{quiz_id}")

    except Exception as e:
        return _render_quiz_error(f"Error generating quiz: {str(e)}")


@app.route("/quiz/view/<int:quiz_id>", methods=["GET"])
@login_required
def quiz_view(quiz_id):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT document, questions_json FROM quizzes WHERE id=?", (quiz_id,))
    row = cur.fetchone()
    conn.close()

    if not row:
        return abort(404)

    document, questions_json = row
    role = normalize_role(session.get("role"))
    if not is_document_accessible(document, role):
        return abort(403)

    questions = json.loads(questions_json)
    return render_template(
        "quiz_view.html",
        quiz_id=quiz_id,
        document=document,
        questions=questions,
        role=role,
    )


@app.route("/quiz/submit/<int:quiz_id>", methods=["POST"])
@login_required
def quiz_submit(quiz_id):
    conn = get_db()
    cur = conn.cursor()
    cur.execute("SELECT document, questions_json FROM quizzes WHERE id=?", (quiz_id,))
    row = cur.fetchone()
    conn.close()

    if not row:
        return abort(404)

    document, questions_json = row
    role = normalize_role(session.get("role"))
    if not is_document_accessible(document, role):
        return abort(403)

    questions = json.loads(questions_json)
    score = 0
    results = []

    for i, q in enumerate(questions):
        user_choice = request.form.get(f"q_{i}", "").strip().upper()
        correct_choice = q.get("answer", "").strip().upper()
        is_correct = user_choice == correct_choice
        if is_correct:
            score += 1

        results.append({
            "question": q.get("question"),
            "options": q.get("options"),
            "correct": correct_choice,
            "user": user_choice,
            "is_correct": is_correct,
            "explanation": q.get("explanation"),
        })

    percent = int((score / len(questions)) * 100) if questions else 0
    username = session.get("username", "unknown")
    log_activity(username, "take_quiz", f"Scored {score}/{len(questions)} ({percent}%) on quiz {quiz_id}")

    return render_template(
        "quiz_results.html",
        document=document,
        score=score,
        total=len(questions),
        percent=percent,
        results=results,
        role=role,
    )


# ─────────────────────────────────────────────────────────────────────────────
# ERROR HANDLERS
# ─────────────────────────────────────────────────────────────────────────────
@app.errorhandler(403)
def forbidden(_error):
    return "You do not have permission to access this resource.", 403


@app.errorhandler(413)
def request_entity_too_large(_error):
    return "File too large. Maximum upload size is 32 MB.", 413


# ─────────────────────────────────────────────────────────────────────────────
# STARTUP
# ─────────────────────────────────────────────────────────────────────────────
with app.app_context():
    migrate_to_chroma()

if __name__ == "__main__":
    port = int(os.environ.get("PORT", 5000))
    debug = os.environ.get("FLASK_DEBUG", "true").lower() == "true"
    app.run(host="0.0.0.0", port=port, debug=debug, use_reloader=debug)

